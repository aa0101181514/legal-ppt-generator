"""讀 JSON 大綱，產出法庭風格 .pptx。

JSON 大綱結構範例：

    {
        "meta": {
            "title": "115金重訴1 逃漏稅答辯簡報",
            "case_no": "高雄地院 115 年度金重訴字第 1 號",
            "defendants": ["謝尚人", "謝尚廷", "謝素芬"],
            "charge": "稅捐稽徵法第 41 條",
            "lawyer": "誠遠商務法律事務所",
            "date": "2026-04-17"
        },
        "slides": [
            {"layout": "cover"},
            {
                "layout": "content",
                "section": "壹、犯後態度",
                "title": "補繳稅額與配合調查",
                "main_point": "被告犯後態度良好，已補繳稅款。",
                "bullets": ["已補繳 1.35 億元", "出具行政陳述書"],
                "citations": [],
                "volume_refs": ["115 金重訴 1 號卷第 187 頁（即起訴書第 175 頁）"]
            },
            {
                "layout": "content",
                "section": "參、重複課稅",
                "title": "國稅局計算公式內含重複課稅",
                "main_point": "同一自費收入同時認列薪資與執業所得。",
                "bullets": ["附檔 A：醫師短漏薪資", "附檔 C：推計執業所得"],
                "citations": [
                    {
                        "verified": true,
                        "query": "107 判字 369",
                        "title": "綜合所得稅",
                        "jyear": "107", "jcase": "判", "jno": "369",
                        "quote": "課稅之構成要件事實，多發生於納稅義務人所得支配之範圍…",
                        "search_url": "https://api.dr-lawbot.com/api/search?q=107+判字+369"
                    }
                ],
                "volume_refs": ["115 金重訴 1 號卷第 80 頁（即起訴書第 68 頁）"]
            },
            {"layout": "conclusion", "statements": ["被告無罪", "請求調查 OO"]},
            {"layout": "appendix"}
        ]
    }

使用：
    python generate_pptx.py outline.json output.pptx
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

import ppt_theme as T


# ─── helpers ─────────────────────────────────────────────

def _add_text(tf, text, *, size, color=None, bold=False, align=PP_ALIGN.LEFT,
              font_zh=None, new_para=False):
    """加一段文字到 text frame。第一段用 tf.paragraphs[0]，之後 add_paragraph。"""
    if new_para or tf.paragraphs[0].text != "" or len(tf.paragraphs[0].runs) > 0:
        # 如果第一段已經有內容 → 新增段落
        if new_para:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    T.apply_run_style(run, size=size, color=color, bold=bold, font_zh=font_zh)
    return p


def _new_para(tf, text, *, size, color=None, bold=False, align=PP_ALIGN.LEFT,
              bullet=False, level=0):
    """新增段落並填字。"""
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    run = p.add_run()
    run.text = text
    T.apply_run_style(run, size=size, color=color, bold=bold)
    return p


def _textbox(slide, left, top, width, height, *, fill=None, line=None,
             wrap=True, anchor=MSO_ANCHOR.TOP):
    """建 textbox 並設定邊距/wrap。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    box.line.fill.background() if line is None else None
    return box, tf


def _fill_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─── slide renderers ─────────────────────────────────────

def render_cover(prs, meta):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, T.COLOR_PRIMARY)

    # 主標
    _, tf = _textbox(
        slide,
        T.MARGIN_LEFT, T.SLIDE_HEIGHT * 0.28,
        T.CONTENT_WIDTH, T.SLIDE_HEIGHT * 0.2,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_text(tf, meta.get("title", "法庭答辯簡報"),
              size=T.FONT_SIZE_TITLE_COVER, color=T.COLOR_BG,
              bold=True, align=PP_ALIGN.CENTER)

    # 副資訊
    _, tf2 = _textbox(
        slide,
        T.MARGIN_LEFT, T.SLIDE_HEIGHT * 0.55,
        T.CONTENT_WIDTH, T.SLIDE_HEIGHT * 0.35,
        anchor=MSO_ANCHOR.TOP,
    )
    info_lines = []
    if meta.get("case_no"):
        info_lines.append(f"案號：{meta['case_no']}")
    if meta.get("defendants"):
        defs = meta["defendants"] if isinstance(meta["defendants"], list) else [meta["defendants"]]
        info_lines.append(f"被告：{'、'.join(defs)}")
    if meta.get("charge"):
        info_lines.append(f"起訴罪名：{meta['charge']}")
    if meta.get("lawyer"):
        info_lines.append(f"辯護人：{meta['lawyer']}")
    if meta.get("date"):
        info_lines.append(f"報告日期：{meta['date']}")

    for i, line in enumerate(info_lines):
        if i == 0:
            _add_text(tf2, line, size=T.FONT_SIZE_SUBTITLE_COVER,
                      color=T.COLOR_BG, align=PP_ALIGN.CENTER)
        else:
            _new_para(tf2, line, size=T.FONT_SIZE_SUBTITLE_COVER,
                      color=T.COLOR_BG, align=PP_ALIGN.CENTER)


def render_toc(prs, slides):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, T.COLOR_BG)

    _, tf = _textbox(
        slide, T.MARGIN_LEFT, T.MARGIN_TOP,
        T.CONTENT_WIDTH, T.Inches(0.8) if False else T.SLIDE_HEIGHT * 0.12,
    )
    _add_text(tf, "目次", size=T.FONT_SIZE_TITLE, color=T.COLOR_PRIMARY,
              bold=True, align=PP_ALIGN.LEFT)

    sections = []
    seen = set()
    for s in slides:
        sec = s.get("section")
        if sec and sec not in seen:
            sections.append(sec)
            seen.add(sec)

    _, tf2 = _textbox(
        slide, T.MARGIN_LEFT, T.SLIDE_HEIGHT * 0.22,
        T.CONTENT_WIDTH, T.SLIDE_HEIGHT * 0.7,
    )
    for i, sec in enumerate(sections):
        if i == 0:
            _add_text(tf2, sec, size=T.FONT_SIZE_BODY, color=T.COLOR_BODY)
        else:
            _new_para(tf2, sec, size=T.FONT_SIZE_BODY, color=T.COLOR_BODY)


def render_section(prs, title):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, T.COLOR_SECTION_BG)

    _, tf = _textbox(
        slide, T.MARGIN_LEFT, T.SLIDE_HEIGHT * 0.35,
        T.CONTENT_WIDTH, T.SLIDE_HEIGHT * 0.3,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_text(tf, title, size=T.FONT_SIZE_SECTION_HEADER,
              color=T.COLOR_BG, bold=True, align=PP_ALIGN.CENTER)


def render_content(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, T.COLOR_BG)

    # 標題列（含章節 tag）
    if slide_data.get("section"):
        _, tf_sec = _textbox(
            slide, T.MARGIN_LEFT, T.MARGIN_TOP,
            T.CONTENT_WIDTH, T.Inches(0.4),
        )
        _add_text(tf_sec, slide_data["section"],
                  size=T.FONT_SIZE_SUBTITLE, color=T.COLOR_PRIMARY_LIGHT,
                  bold=True)

    # 大標
    _, tf_title = _textbox(
        slide, T.MARGIN_LEFT, T.MARGIN_TOP + T.Inches(0.45),
        T.CONTENT_WIDTH, T.Inches(0.7),
    )
    _add_text(tf_title, slide_data.get("title", ""),
              size=T.FONT_SIZE_TITLE, color=T.COLOR_PRIMARY, bold=True)

    # 底下分隔線
    line_top = T.MARGIN_TOP + T.Inches(1.2)
    line = slide.shapes.add_connector(
        1,  # straight line
        T.MARGIN_LEFT, line_top,
        T.MARGIN_LEFT + T.CONTENT_WIDTH, line_top,
    )
    line.line.color.rgb = T.COLOR_LINE

    # 核心主張（如果有）
    body_top = line_top + T.Inches(0.2)
    if slide_data.get("main_point"):
        _, tf_mp = _textbox(
            slide, T.MARGIN_LEFT, body_top,
            T.CONTENT_WIDTH, T.Inches(0.6),
        )
        _add_text(tf_mp, "● 核心主張：",
                  size=T.FONT_SIZE_BODY, color=T.COLOR_ACCENT, bold=True)
        run = tf_mp.paragraphs[0].add_run()
        run.text = slide_data["main_point"]
        T.apply_run_style(run, size=T.FONT_SIZE_BODY, color=T.COLOR_BODY)
        body_top += T.Inches(0.7)

    # Bullets
    bullets = slide_data.get("bullets", [])
    if bullets:
        _, tf_b = _textbox(
            slide, T.MARGIN_LEFT, body_top,
            T.CONTENT_WIDTH, T.Inches(2.5),
        )
        for i, b in enumerate(bullets):
            if i == 0:
                _add_text(tf_b, f"・{b}",
                          size=T.FONT_SIZE_BODY, color=T.COLOR_BODY)
            else:
                _new_para(tf_b, f"・{b}",
                          size=T.FONT_SIZE_BODY, color=T.COLOR_BODY)
        body_top += T.Inches(0.5 * max(1, len(bullets)) + 0.3)

    # 判決引用
    citations = slide_data.get("citations", [])
    if citations:
        for cite in citations:
            body_top += T.Inches(0.1)
            box_h = T.Inches(1.6)
            box, tf_c = _textbox(
                slide, T.MARGIN_LEFT, body_top,
                T.CONTENT_WIDTH, box_h,
                fill=T.COLOR_QUOTE_BG,
            )
            tf_c.margin_left = tf_c.margin_right = T.Inches(0.2)
            tf_c.margin_top = tf_c.margin_bottom = T.Inches(0.1)

            # 判決字號 + 驗證標籤
            verified = cite.get("verified", False)
            mark = T.VERIFIED_MARK if verified else T.UNVERIFIED_MARK
            mark_color = T.COLOR_PRIMARY if verified else T.COLOR_WARNING
            case_line = _format_citation_line(cite)

            _add_text(tf_c, case_line,
                      size=T.FONT_SIZE_BODY_SUB,
                      color=T.COLOR_PRIMARY, bold=True)
            run_mark = tf_c.paragraphs[0].add_run()
            run_mark.text = f"  {mark}"
            T.apply_run_style(run_mark, size=T.FONT_SIZE_FOOTER, color=mark_color)

            # 引文
            quote = cite.get("quote") or cite.get("main_preview") or ""
            if quote:
                quote = quote.strip()
                if len(quote) > 150:
                    quote = quote[:150].rstrip() + "……"
                _new_para(tf_c, f"「{quote}」",
                          size=T.FONT_SIZE_QUOTE, color=T.COLOR_QUOTE_TEXT)

            body_top += box_h + T.Inches(0.1)

    # 卷頁引用（footer）
    vol_refs = slide_data.get("volume_refs", [])
    if vol_refs:
        footer_top = T.SLIDE_HEIGHT - T.MARGIN_BOTTOM - T.Inches(0.5)
        _, tf_f = _textbox(
            slide, T.MARGIN_LEFT, footer_top,
            T.CONTENT_WIDTH, T.Inches(0.4),
        )
        ref_str = "　".join(vol_refs)
        _add_text(tf_f, f"卷證出處：{ref_str}",
                  size=T.FONT_SIZE_FOOTER, color=T.COLOR_PRIMARY_LIGHT)


def _format_citation_line(cite):
    """把判決資料格式化成一行引用標題。"""
    jyear = cite.get("jyear")
    jcase = cite.get("jcase")
    jno = cite.get("jno")
    title = cite.get("title") or ""

    # 從 doc_id 推法院名
    doc_id = cite.get("doc_id") or ""
    prefix = doc_id.split(",")[0] if "," in doc_id else ""
    court_map = {
        "TPSM": "最高法院", "TPSV": "最高法院",
        "TPAA": "最高行政法院",
        "TPHV": "台灣高等法院", "TCHV": "台中高分院",
        "TNHV": "台南高分院", "KSHV": "高雄高分院",
        "TPBA": "台北高等行政法院", "TCBA": "台中高等行政法院",
        "KSBA": "高雄高等行政法院",
    }
    court = court_map.get(prefix, "")

    if jyear and jcase and jno:
        line = f"{court} {jyear} 年度{jcase}字第 {jno} 號判決"
    else:
        line = cite.get("query", "（未指定字號）")
    if title:
        line += f"（{title}）"
    return line.strip()


def render_conclusion(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, T.COLOR_BG)

    _, tf = _textbox(
        slide, T.MARGIN_LEFT, T.MARGIN_TOP,
        T.CONTENT_WIDTH, T.Inches(0.8),
    )
    _add_text(tf, "結論與聲明",
              size=T.FONT_SIZE_TITLE, color=T.COLOR_PRIMARY, bold=True)

    statements = slide_data.get("statements", [])
    _, tf_s = _textbox(
        slide, T.MARGIN_LEFT, T.MARGIN_TOP + T.Inches(1.2),
        T.CONTENT_WIDTH, T.Inches(5.0),
    )
    for i, s in enumerate(statements):
        text = f"{i+1}. {s}"
        if i == 0:
            _add_text(tf_s, text, size=T.FONT_SIZE_BODY, color=T.COLOR_BODY)
        else:
            _new_para(tf_s, text, size=T.FONT_SIZE_BODY, color=T.COLOR_BODY)


def render_appendix(prs, citations_flat):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _fill_background(slide, T.COLOR_BG)

    _, tf = _textbox(
        slide, T.MARGIN_LEFT, T.MARGIN_TOP,
        T.CONTENT_WIDTH, T.Inches(0.8),
    )
    _add_text(tf, "附錄：本簡報引用判決清單（均經 dr-lawbot 驗證）",
              size=T.FONT_SIZE_TITLE, color=T.COLOR_PRIMARY, bold=True)

    _, tf_list = _textbox(
        slide, T.MARGIN_LEFT, T.MARGIN_TOP + T.Inches(1.2),
        T.CONTENT_WIDTH, T.Inches(5.5),
    )

    for i, c in enumerate(citations_flat):
        line = f"{i+1}. {_format_citation_line(c)}"
        if c.get("search_url"):
            line += f"\n    {c['search_url']}"
        if i == 0:
            _add_text(tf_list, line, size=T.FONT_SIZE_BODY_SUB, color=T.COLOR_BODY)
        else:
            _new_para(tf_list, line, size=T.FONT_SIZE_BODY_SUB, color=T.COLOR_BODY)


# ─── main ───────────────────────────────────────────────

def generate(outline: dict, output_path: str | Path) -> Path:
    prs = Presentation()
    prs.slide_width = T.SLIDE_WIDTH
    prs.slide_height = T.SLIDE_HEIGHT

    slides = outline.get("slides", [])
    meta = outline.get("meta", {})

    # 統計全部 citations 供 appendix 用
    all_citations = []
    sections_seen = set()

    for s in slides:
        layout = s.get("layout", "content")
        if layout == "cover":
            render_cover(prs, meta)
        elif layout == "toc":
            render_toc(prs, slides)
        elif layout == "section":
            render_section(prs, s.get("title", ""))
        elif layout == "content":
            # 若遇到新章節 → 先插入 section 分隔頁
            sec = s.get("section")
            if sec and sec not in sections_seen:
                sections_seen.add(sec)
                render_section(prs, sec)
            render_content(prs, s)
            for c in s.get("citations", []):
                all_citations.append(c)
        elif layout == "conclusion":
            render_conclusion(prs, s)
        elif layout == "appendix":
            render_appendix(prs, all_citations)
        else:
            # 未知 layout fallback 為 content
            render_content(prs, s)

    output_path = Path(output_path)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法：python generate_pptx.py <outline.json> <output.pptx>")
        sys.exit(1)
    json_path = Path(sys.argv[1])
    out = Path(sys.argv[2])
    outline = json.loads(json_path.read_text(encoding="utf-8"))
    result = generate(outline, out)
    print(f"✓ 已產出：{result.resolve()}")
